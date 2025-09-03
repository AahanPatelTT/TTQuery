#!/usr/bin/env python3
"""
Image extraction module for Synapse

This module extracts images/diagrams from PDFs and other documents, saving them as 
independent files with rich metadata for better retrieval and contextual understanding.

Key features:
- Extract images from PDFs with page context
- Generate descriptive filenames based on document and page context
- Store metadata about image source, context, and technical content
- Support for various image formats and OCR enhancement
- Integration with existing Synapse parsing pipeline
"""

import hashlib
import json
import logging
import os
import re
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any, Set
import base64
import asyncio
import concurrent.futures
from functools import partial
try:
    import imagehash
    IMAGEHASH_AVAILABLE = True
except ImportError:
    IMAGEHASH_AVAILABLE = False
    logging.warning("imagehash library not available. Install with: pip install imagehash")


@dataclass
class ExtractedImage:
    """Metadata for an extracted image."""
    image_id: str                    # Unique identifier
    source_document: str             # Original document path
    source_type: str                 # pdf, pptx, etc.
    page_number: Optional[int]       # Page/slide where image was found
    image_path: str                  # Path to extracted image file
    image_format: str                # png, jpg, etc.
    width: int                       # Image width in pixels
    height: int                      # Image height in pixels
    file_size: int                   # Image file size in bytes
    extraction_method: str           # Method used to extract (pymupdf, unstructured, etc.)
    document_context: str            # Surrounding text from the same page
    ocr_text: str                    # OCR text from the image
    caption_text: str                # AI-generated caption
    technical_keywords: List[str]    # Extracted technical terms
    image_type: str                  # diagram, chart, photo, etc.
    content_hash: Optional[str] = None  # Content-based hash for deduplication
    is_duplicate: bool = False       # Whether this image is a duplicate of another
    original_image_id: Optional[str] = None  # ID of the original if this is a duplicate
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization."""
        return asdict(self)


class ImageDeduplicationCache:
    """Cache for tracking duplicate images during extraction."""
    
    def __init__(self, similarity_threshold: int = 10):
        self.content_hash_to_image: Dict[str, ExtractedImage] = {}
        self.processed_hashes: Set[str] = set()
        # For perceptual hashing - maps hash to ExtractedImage
        self.perceptual_hashes: Dict[str, ExtractedImage] = {}
        self.similarity_threshold = similarity_threshold  # Hamming distance threshold for similar images
    
    def is_duplicate(self, content_hash: str, perceptual_hash: Optional[str] = None) -> Tuple[bool, Optional[str]]:
        """Check if an image with this content hash or similar perceptual hash has already been processed.
        
        Returns:
            Tuple[bool, Optional[str]]: (is_duplicate, original_hash_key)
        """
        # First check exact content match
        if content_hash in self.processed_hashes:
            return True, content_hash
        
        # If perceptual hashing is available and we have a perceptual hash, check for similar images
        if IMAGEHASH_AVAILABLE and perceptual_hash:
            try:
                current_hash = imagehash.hex_to_hash(perceptual_hash)
                for existing_hash_str, existing_image in self.perceptual_hashes.items():
                    existing_hash = imagehash.hex_to_hash(existing_hash_str)
                    # Calculate Hamming distance
                    distance = current_hash - existing_hash
                    if distance <= self.similarity_threshold:
                        logging.debug(f"Found similar image with distance {distance} (threshold: {self.similarity_threshold})")
                        return True, existing_image.content_hash
            except Exception as e:
                logging.debug(f"Error comparing perceptual hashes: {e}")
        
        return False, None
    
    def add_image(self, image: ExtractedImage, content_hash: str, perceptual_hash: Optional[str] = None) -> None:
        """Add an image to the cache."""
        self.content_hash_to_image[content_hash] = image
        self.processed_hashes.add(content_hash)
        image.content_hash = content_hash
        
        # Also store perceptual hash if available
        if IMAGEHASH_AVAILABLE and perceptual_hash:
            self.perceptual_hashes[perceptual_hash] = image
    
    def get_original_image(self, content_hash: str) -> Optional[ExtractedImage]:
        """Get the original image for a given content hash."""
        return self.content_hash_to_image.get(content_hash)
    
    def create_duplicate_reference(self, original_image: ExtractedImage, 
                                 source_document: str, page_number: Optional[int],
                                 document_context: str, source_type: str) -> ExtractedImage:
        """Create a duplicate reference that points to the original image."""
        duplicate_id = self._generate_duplicate_id(source_document, page_number, original_image.image_id)
        
        return ExtractedImage(
            image_id=duplicate_id,
            source_document=source_document,
            source_type=source_type,
            page_number=page_number,
            image_path=original_image.image_path,  # Point to the same file
            image_format=original_image.image_format,
            width=original_image.width,
            height=original_image.height,
            file_size=original_image.file_size,
            extraction_method=original_image.extraction_method,
            document_context=document_context,
            ocr_text=original_image.ocr_text,  # Reuse OCR text
            caption_text=original_image.caption_text,  # Reuse caption
            technical_keywords=original_image.technical_keywords,
            image_type=original_image.image_type,
            content_hash=original_image.content_hash,
            is_duplicate=True,
            original_image_id=original_image.image_id
        )
    
    def _generate_duplicate_id(self, source_document: str, page_number: Optional[int], original_id: str) -> str:
        """Generate a unique ID for a duplicate image reference."""
        source_str = f"duplicate:{source_document}:page_{page_number}:original_{original_id}"
        return hashlib.sha256(source_str.encode()).hexdigest()


def optimize_image_for_ocr(image_data: bytes, max_dimension: int = 2048, fast_mode: bool = False) -> bytes:
    """Optimize large images by reducing resolution while preserving OCR quality.
    
    Args:
        image_data: Original image data
        max_dimension: Maximum width or height for the optimized image
    
    Returns:
        Optimized image data (or original if no optimization needed)
    """
    try:
        from PIL import Image
        import io
        
        with Image.open(io.BytesIO(image_data)) as img:
            width, height = img.size
            
            # In fast mode, use more aggressive optimization
            if fast_mode:
                max_dimension = min(max_dimension, 1024)  # Smaller max size in fast mode
            
            # Only optimize if image is larger than max_dimension
            if max(width, height) <= max_dimension:
                return image_data
            
            # Calculate new dimensions maintaining aspect ratio
            if width > height:
                new_width = max_dimension
                new_height = int(height * (max_dimension / width))
            else:
                new_height = max_dimension
                new_width = int(width * (max_dimension / height))
            
            # Choose resampling method based on mode
            if fast_mode:
                # Faster but lower quality resampling
                img_resized = img.resize((new_width, new_height), Image.Resampling.BILINEAR)
            else:
                # Higher quality but slower resampling
                img_resized = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
            
            # Save optimized image
            img_buffer = io.BytesIO()
            if fast_mode:
                img_resized.save(img_buffer, format='PNG', optimize=False)  # Skip optimization in fast mode
            else:
                img_resized.save(img_buffer, format='PNG', optimize=True)
            optimized_data = img_buffer.getvalue()
            
            logging.debug(f"Optimized image from {width}x{height} to {new_width}x{new_height} (size: {len(image_data)} -> {len(optimized_data)} bytes)")
            return optimized_data
            
    except Exception as e:
        logging.warning(f"Failed to optimize image, using original: {e}")
        return image_data


def compute_image_hashes(image_data: bytes) -> Tuple[str, Optional[str]]:
    """Compute both content-based hash and perceptual hash for an image.
    
    Returns:
        Tuple[str, Optional[str]]: (content_hash, perceptual_hash)
    """
    try:
        from PIL import Image
        import io
        
        # Open the image
        with Image.open(io.BytesIO(image_data)) as img:
            # Convert to RGB to normalize color spaces
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            # Compute perceptual hash if available
            perceptual_hash = None
            if IMAGEHASH_AVAILABLE:
                try:
                    # Use average hash for good balance of speed and accuracy
                    # You can also use phash (perceptual hash) for better accuracy but slower performance
                    hash_obj = imagehash.average_hash(img, hash_size=16)
                    perceptual_hash = str(hash_obj)
                except Exception as e:
                    logging.debug(f"Failed to compute perceptual hash: {e}")
            
            # Resize to a standard size for content comparison (handles minor scaling differences)
            img_normalized = img.resize((512, 512), Image.Resampling.LANCZOS)
            
            # Convert back to bytes for content hash
            img_buffer = io.BytesIO()
            img_normalized.save(img_buffer, format='PNG')
            normalized_data = img_buffer.getvalue()
            
            # Hash the normalized image data
            content_hash = hashlib.sha256(normalized_data).hexdigest()
            
            return content_hash, perceptual_hash
            
    except Exception as e:
        logging.warning(f"Failed to compute image hashes, using raw data hash: {e}")
        # Fallback to raw data hash
        content_hash = hashlib.sha256(image_data).hexdigest()
        return content_hash, None


def compute_image_content_hash(image_data: bytes) -> str:
    """Compute a content-based hash for an image to detect duplicates.
    
    This function is kept for backward compatibility.
    """
    content_hash, _ = compute_image_hashes(image_data)
    return content_hash


def compute_image_file_hash(image_path: str) -> str:
    """Compute a content-based hash for an image file."""
    try:
        with open(image_path, 'rb') as f:
            image_data = f.read()
        content_hash, _ = compute_image_hashes(image_data)
        return content_hash
    except Exception as e:
        logging.error(f"Failed to compute hash for image file {image_path}: {e}")
        return ""


class PDFImageExtractor:
    """Extract images from PDF documents."""
    
    def __init__(self, output_dir: str = "artifacts/extracted_images", enable_captioning: bool = True, 
                 dedup_cache: Optional[ImageDeduplicationCache] = None, fast_mode: bool = False): 
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.enable_captioning = enable_captioning and not fast_mode  # Disable captioning in fast mode
        self.dedup_cache = dedup_cache or ImageDeduplicationCache()
        self.fast_mode = fast_mode
        # Cache for fast OCR processing
        self._fast_ocr_cache = {}
        # Thread pool for async processing
        self._executor = None
        
    def extract_images_pymupdf(self, pdf_path: str, page_context: Optional[Dict[int, str]] = None) -> List[ExtractedImage]:
        """Extract images using PyMuPDF (fast, good for most PDFs)."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            logging.warning("PyMuPDF not available for image extraction")
            return []
        
        extracted_images = []
        
        try:
            doc = fitz.open(pdf_path)
            doc_name = Path(pdf_path).stem
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                image_list = page.get_images()
                
                # Get page text for context
                page_text = page_context.get(page_num + 1, "") if page_context else page.get_text()
                
                for img_index, img in enumerate(image_list):
                    try:
                        # Get image data
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        
                        # Skip if image is too small (likely not a meaningful diagram)
                        if pix.width < 100 or pix.height < 100:
                            pix = None
                            continue
                        
                        # Get image data for deduplication check
                        if pix.n - pix.alpha < 4:  # GRAY or RGB
                            image_data = pix.tobytes("png")
                        else:  # CMYK: convert to RGB first
                            pix1 = fitz.Pixmap(fitz.csRGB, pix)
                            image_data = pix1.tobytes("png")
                            pix1 = None
                        
                        # Optimize image for better processing (especially for large images)
                        optimized_image_data = optimize_image_for_ocr(image_data, fast_mode=self.fast_mode)
                        
                        # Compute both content and perceptual hashes for deduplication
                        content_hash, perceptual_hash = compute_image_hashes(image_data)
                        
                        # Check if this image is a duplicate
                        is_dup, original_hash = self.dedup_cache.is_duplicate(content_hash, perceptual_hash)
                        if is_dup:
                            original_image = self.dedup_cache.get_original_image(original_hash)
                            if original_image:
                                logging.debug(f"Found duplicate image in {pdf_path} page {page_num + 1}, referencing original (hash: {original_hash[:8]})")
                                # Create a duplicate reference instead of saving the image again
                                duplicate_ref = self.dedup_cache.create_duplicate_reference(
                                    original_image, 
                                    os.path.abspath(pdf_path),
                                    page_num + 1,
                                    page_text[:2000],
                                    "pdf"
                                )
                                extracted_images.append(duplicate_ref)
                                pix = None
                                continue
                        
                        # This is a new unique image - save it
                        image_id = self._generate_image_id(pdf_path, page_num + 1, img_index)
                        image_filename = f"{doc_name}_p{page_num + 1:03d}_img{img_index:02d}_{image_id[:8]}.png"
                        image_path = self.output_dir / image_filename
                        
                        # Save the optimized image data
                        with open(image_path, 'wb') as f:
                            f.write(optimized_image_data)
                        
                        # Extract text using OCR if available
                        ocr_text = self._extract_ocr_text(str(image_path))
                        
                        # Generate caption
                        caption_text = self._generate_caption(str(image_path))
                        
                        # Detect image type and extract keywords
                        image_type, keywords = self._analyze_image_content(ocr_text, caption_text, page_text)
                        
                        # Create metadata
                        extracted_image = ExtractedImage(
                            image_id=image_id,
                            source_document=os.path.abspath(pdf_path),
                            source_type="pdf",
                            page_number=page_num + 1,
                            image_path=str(image_path.absolute()),
                            image_format="png",
                            width=pix.width,
                            height=pix.height,
                            file_size=os.path.getsize(image_path),
                            extraction_method="pymupdf",
                            document_context=page_text[:2000],  # Limit context size
                            ocr_text=ocr_text,
                            caption_text=caption_text,
                            technical_keywords=keywords,
                            image_type=image_type
                        )
                        
                        # Add to deduplication cache
                        self.dedup_cache.add_image(extracted_image, content_hash, perceptual_hash)
                        
                        extracted_images.append(extracted_image)
                        logging.debug(f"Extracted unique image: {image_filename}")
                        
                        pix = None
                        
                    except Exception as e:
                        logging.warning(f"Failed to extract image {img_index} from page {page_num + 1}: {e}")
                        continue
                        
            doc.close()
            
        except Exception as e:
            logging.error(f"Failed to extract images from {pdf_path}: {e}")
            
        return extracted_images
    
    def _get_executor(self):
        """Get or create thread pool executor for async processing."""
        if self._executor is None:
            self._executor = concurrent.futures.ThreadPoolExecutor(max_workers=4)
        return self._executor
    
    async def extract_images_async(self, pdf_path: str, page_context: Optional[Dict[int, str]] = None) -> List[ExtractedImage]:
        """Async version of image extraction for real-time processing."""
        loop = asyncio.get_event_loop()
        executor = self._get_executor()
        
        # Run the synchronous extraction in a thread pool
        return await loop.run_in_executor(
            executor, 
            partial(self.extract_images_pymupdf, pdf_path, page_context)
        )
    
    def extract_images_unstructured(self, pdf_path: str) -> List[ExtractedImage]:
        """Extract images using unstructured.io (more comprehensive but slower)."""
        try:
            from unstructured.partition.pdf import partition_pdf
        except ImportError:
            logging.warning("unstructured library not available for image extraction")
            return []
        
        extracted_images = []
        
        try:
            # Use unstructured to get elements with image detection
            elements = partition_pdf(
                filename=pdf_path,
                strategy="hi_res",
                extract_images_in_pdf=True,
                infer_table_structure=True
            )
            
            doc_name = Path(pdf_path).stem
            
            for element in elements:
                if hasattr(element, 'metadata') and hasattr(element.metadata, 'image_path'):
                    try:
                        source_image_path = element.metadata.image_path
                        
                        if not os.path.exists(source_image_path):
                            continue
                        
                        # Get image info
                        from PIL import Image
                        with Image.open(source_image_path) as img:
                            width, height = img.size
                            img_format = img.format or "PNG"
                        
                        # Generate unique filename and copy to our output directory
                        page_num = getattr(element.metadata, 'page_number', 0)
                        image_id = self._generate_image_id(pdf_path, page_num, 0)
                        image_filename = f"{doc_name}_p{page_num:03d}_unstr_{image_id[:8]}.{img_format.lower()}"
                        image_path = self.output_dir / image_filename
                        
                        # Copy image to our directory
                        import shutil
                        shutil.copy2(source_image_path, image_path)
                        
                        # Extract context from nearby elements
                        page_context = self._get_page_context_unstructured(elements, page_num)
                        
                        # Extract text using OCR
                        ocr_text = self._extract_ocr_text(str(image_path))
                        
                        # Generate caption
                        caption_text = self._generate_caption(str(image_path))
                        
                        # Analyze content
                        image_type, keywords = self._analyze_image_content(ocr_text, caption_text, page_context)
                        
                        extracted_image = ExtractedImage(
                            image_id=image_id,
                            source_document=os.path.abspath(pdf_path),
                            source_type="pdf",
                            page_number=page_num,
                            image_path=str(image_path.absolute()),
                            image_format=img_format.lower(),
                            width=width,
                            height=height,
                            file_size=os.path.getsize(image_path),
                            extraction_method="unstructured",
                            document_context=page_context,
                            ocr_text=ocr_text,
                            caption_text=caption_text,
                            technical_keywords=keywords,
                            image_type=image_type
                        )
                        
                        extracted_images.append(extracted_image)
                        logging.debug(f"Extracted image via unstructured: {image_filename}")
                        
                    except Exception as e:
                        logging.warning(f"Failed to process unstructured image: {e}")
                        continue
                        
        except Exception as e:
            logging.error(f"Failed to extract images via unstructured from {pdf_path}: {e}")
            
        return extracted_images
    
    def _generate_image_id(self, source_path: str, page_num: int, img_index: int) -> str:
        """Generate a unique ID for an extracted image."""
        source_str = f"{source_path}:page_{page_num}:img_{img_index}"
        return hashlib.sha256(source_str.encode()).hexdigest()
    
    def _extract_ocr_text(self, image_path: str) -> str:
        """Extract text from image using OCR."""
        try:
            import pytesseract
            from PIL import Image
            
            # In fast mode, use simpler OCR configuration
            if self.fast_mode:
                with Image.open(image_path) as img:
                    # Use faster, simpler OCR config
                    ocr_text = pytesseract.image_to_string(img, config='--psm 6')
                    return self._correct_technical_ocr_errors(ocr_text.strip())
            
            # Standard mode with enhanced configuration
            with Image.open(image_path) as img:
                # Use enhanced OCR configuration for technical diagrams
                custom_config = r'--oem 3 --psm 6 -c tessedit_char_whitelist=0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz ×x-/()[]{}:.,;'
                ocr_text = pytesseract.image_to_string(img, config=custom_config)
                
                # Apply common technical OCR corrections
                ocr_text = self._correct_technical_ocr_errors(ocr_text)
                
                return ocr_text.strip()
                
        except Exception as e:
            logging.debug(f"OCR extraction failed for {image_path}: {e}")
            return ""
    
    def _generate_caption(self, image_path: str) -> str:
        """Generate an AI caption for the image."""
        # Skip captioning if disabled
        if not self.enable_captioning:
            return ""
            
        try:
            # Cache the captioner pipeline to avoid reloading for each image
            if not hasattr(self, '_captioner'):
                try:
                    from transformers import pipeline
                    from PIL import Image
                    
                    logging.info("Loading BLIP image captioning model (this may take a moment)...")
                    self._captioner = pipeline("image-to-text", model="Salesforce/blip-image-captioning-base")
                    logging.info("BLIP model loaded successfully")
                except Exception as model_error:
                    logging.warning(f"Failed to load BLIP model (likely due to rate limits): {model_error}")
                    # Fallback: disable captioning for this session
                    self._captioner = None
                    return ""
            
            if self._captioner is None:
                return ""
                
            from PIL import Image
            with Image.open(image_path) as img:
                result = self._captioner(img, max_new_tokens=64)
                if isinstance(result, list) and result:
                    return result[0].get("generated_text", "").strip()
                    
        except Exception as e:
            logging.debug(f"Image captioning failed for {image_path}: {e}")
            
        return ""
    
    def _correct_technical_ocr_errors(self, ocr_text: str) -> str:
        """Apply corrections for common OCR errors in technical diagrams."""
        if not ocr_text:
            return ocr_text
        
        corrections = {
            # Common OCR errors in technical diagrams
            "12C": "I2C",
            "13C": "I2C", 
            "SPI ": "SPI ",
            "UARTx": "UART x",
            "GPIOx": "GPIO x",
            "CANFD": "CAN-FD",
            "CANFDx": "CAN-FD x",
            "SPIx": "SPI x",
            "I2Cx": "I2C x",
            " x ": " x",  # normalize spacing around x
            "×": "x",      # replace multiplication symbol with x
            "Timer": "Timer",
            "WDT": "WDT",
            "RAS": "RAS",
            "Cluster": "Cluster",
            "CPU": "CPU",
            "Core": "Core",
            "Cache": "Cache",
            "L1": "L1",
            "L2": "L2",
            "L3": "L3"
        }
        
        result = ocr_text
        for error, correction in corrections.items():
            result = result.replace(error, correction)
        
        return result
    
    def _analyze_image_content(self, ocr_text: str, caption_text: str, page_context: str) -> Tuple[str, List[str]]:
        """Analyze image content to determine type and extract technical keywords."""
        
        # Combine all text sources for analysis
        all_text = f"{ocr_text} {caption_text} {page_context}".lower()
        
        # Determine image type based on content
        image_type = "diagram"  # default
        
        if any(word in all_text for word in ["chart", "graph", "plot", "performance", "benchmark"]):
            image_type = "chart"
        elif any(word in all_text for word in ["block", "architecture", "system", "connection", "interface"]):
            image_type = "block_diagram"
        elif any(word in all_text for word in ["flow", "process", "pipeline", "state"]):
            image_type = "flowchart"
        elif any(word in all_text for word in ["table", "matrix", "grid"]):
            image_type = "table"
        elif any(word in all_text for word in ["photo", "image", "picture"]):
            image_type = "photo"
        
        # Extract technical keywords
        technical_patterns = [
            r'\b(?:CPU|GPU|TPU|DSP|RISC-V|ARM|x86)\b',
            r'\b(?:I2C|SPI|UART|USB|PCIe|CAN|Ethernet)\b',
            r'\b(?:L1|L2|L3)\s*(?:cache|Cache)\b',
            r'\b(?:Timer|WDT|RAS|Cluster|Core|Cache)\b',
            r'\b(?:MHz|GHz|MB|GB|TB|Mbps|Gbps)\b',
            r'\b(?:interrupt|vector|peripheral|interface)\b',
            r'\b(?:Ascalon|Tensix|Alexandria|Blackhole)\b',  # TT-specific terms
            r'\b\d+(?:x\d+|\s*x\s*\d+)\b',  # Dimensions like "2x4", "8 x 8"
        ]
        
        keywords = []
        for pattern in technical_patterns:
            matches = re.findall(pattern, all_text, re.IGNORECASE)
            keywords.extend([match.strip() for match in matches])
        
        # Remove duplicates and normalize
        keywords = list(set([kw.strip() for kw in keywords if len(kw.strip()) > 1]))
        
        return image_type, keywords
    
    def _get_page_context_unstructured(self, elements: List, target_page: int) -> str:
        """Extract context text from the same page using unstructured elements."""
        page_texts = []
        
        for element in elements:
            if hasattr(element, 'metadata') and hasattr(element.metadata, 'page_number'):
                if element.metadata.page_number == target_page:
                    if hasattr(element, 'text') and element.text:
                        page_texts.append(element.text.strip())
        
        return " ".join(page_texts)[:2000]  # Limit context size


class PPTXImageExtractor:
    """Extract images from PowerPoint presentations."""
    
    def __init__(self, output_dir: str = "artifacts/extracted_images", enable_captioning: bool = True,
                 dedup_cache: Optional[ImageDeduplicationCache] = None, fast_mode: bool = False):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.enable_captioning = enable_captioning and not fast_mode  # Disable captioning in fast mode
        self.dedup_cache = dedup_cache or ImageDeduplicationCache()
        self.fast_mode = fast_mode
    
    def extract_images(self, pptx_path: str) -> List[ExtractedImage]:
        """Extract images from PPTX file."""
        try:
            from pptx import Presentation
            from PIL import Image
        except ImportError:
            logging.warning("python-pptx or PIL not available for PPTX image extraction")
            return []
        
        extracted_images = []
        
        try:
            pres = Presentation(pptx_path)
            doc_name = Path(pptx_path).stem
            
            for slide_idx, slide in enumerate(pres.slides):
                # Get slide text for context
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_texts.append(shape.text.strip())
                slide_context = " ".join(slide_texts)
                
                # Extract images from slide
                img_idx = 0
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "image"):
                            # Extract embedded image
                            image_stream = shape.image.blob
                            
                            # Optimize image for better processing
                            optimized_image_data = optimize_image_for_ocr(image_stream, fast_mode=self.fast_mode)
                            
                            # Compute both content and perceptual hashes for deduplication
                            content_hash, perceptual_hash = compute_image_hashes(image_stream)
                            
                            # Check if this image is a duplicate
                            is_dup, original_hash = self.dedup_cache.is_duplicate(content_hash, perceptual_hash)
                            if is_dup:
                                original_image = self.dedup_cache.get_original_image(original_hash)
                                if original_image:
                                    logging.debug(f"Found duplicate image in {pptx_path} slide {slide_idx + 1}, referencing original (hash: {original_hash[:8]})")
                                    # Create a duplicate reference instead of saving the image again
                                    duplicate_ref = self.dedup_cache.create_duplicate_reference(
                                        original_image,
                                        os.path.abspath(pptx_path),
                                        slide_idx + 1,
                                        slide_context[:2000],
                                        "pptx"
                                    )
                                    extracted_images.append(duplicate_ref)
                                    img_idx += 1
                                    continue
                            
                            # This is a new unique image - save it
                            image_id = self._generate_image_id(pptx_path, slide_idx + 1, img_idx)
                            image_filename = f"{doc_name}_s{slide_idx + 1:03d}_img{img_idx:02d}_{image_id[:8]}.png"
                            image_path = self.output_dir / image_filename
                            
                            # Save optimized image
                            with open(image_path, 'wb') as f:
                                f.write(optimized_image_data)
                            
                            # Get image dimensions
                            with Image.open(image_path) as img:
                                width, height = img.size
                            
                            # Process image content
                            ocr_text = self._extract_ocr_text(str(image_path))
                            caption_text = self._generate_caption(str(image_path))
                            image_type, keywords = self._analyze_image_content(ocr_text, caption_text, slide_context)
                            
                            extracted_image = ExtractedImage(
                                image_id=image_id,
                                source_document=os.path.abspath(pptx_path),
                                source_type="pptx",
                                page_number=slide_idx + 1,
                                image_path=str(image_path.absolute()),
                                image_format="png",
                                width=width,
                                height=height,
                                file_size=os.path.getsize(image_path),
                                extraction_method="python-pptx",
                                document_context=slide_context[:2000],
                                ocr_text=ocr_text,
                                caption_text=caption_text,
                                technical_keywords=keywords,
                                image_type=image_type
                            )
                            
                            # Add to deduplication cache
                            self.dedup_cache.add_image(extracted_image, content_hash, perceptual_hash)
                            
                            extracted_images.append(extracted_image)
                            logging.debug(f"Extracted unique PPTX image: {image_filename}")
                            
                            img_idx += 1
                            
                    except Exception as e:
                        logging.warning(f"Failed to extract image from slide {slide_idx + 1}: {e}")
                        continue
                        
        except Exception as e:
            logging.error(f"Failed to extract images from {pptx_path}: {e}")
            
        return extracted_images
    
    def _generate_image_id(self, source_path: str, slide_num: int, img_index: int) -> str:
        """Generate a unique ID for an extracted image."""
        source_str = f"{source_path}:slide_{slide_num}:img_{img_index}"
        return hashlib.sha256(source_str.encode()).hexdigest()
    
    def _extract_ocr_text(self, image_path: str) -> str:
        """Extract text from image using OCR."""
        try:
            import pytesseract
            from PIL import Image
            
            # In fast mode, use simpler OCR configuration
            if self.fast_mode:
                with Image.open(image_path) as img:
                    # Use faster, simpler OCR config
                    ocr_text = pytesseract.image_to_string(img, config='--psm 6')
                    return self._correct_technical_ocr_errors(ocr_text.strip())
            
            # Standard mode with enhanced configuration
            with Image.open(image_path) as img:
                custom_config = r'--oem 3 --psm 6 -c tessedit_char_whitelist=0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz ×x-/()[]{}:.,;'
                ocr_text = pytesseract.image_to_string(img, config=custom_config)
                return self._correct_technical_ocr_errors(ocr_text.strip())
                
        except Exception as e:
            logging.debug(f"OCR extraction failed for {image_path}: {e}")
            return ""
    
    def _generate_caption(self, image_path: str) -> str:
        """Generate an AI caption for the image."""
        try:
            from transformers import pipeline
            from PIL import Image
            
            captioner = pipeline("image-to-text", model="Salesforce/blip-image-captioning-base")
            
            with Image.open(image_path) as img:
                result = captioner(img, max_new_tokens=64)
                if isinstance(result, list) and result:
                    return result[0].get("generated_text", "").strip()
                    
        except Exception as e:
            logging.debug(f"Image captioning failed for {image_path}: {e}")
            
        return ""
    
    def _correct_technical_ocr_errors(self, ocr_text: str) -> str:
        """Apply corrections for common OCR errors in technical diagrams."""
        if not ocr_text:
            return ocr_text
        
        corrections = {
            "12C": "I2C", "13C": "I2C", "RAS": "RAS", "Cluster": "Cluster",
            "CPU": "CPU", "Core": "Core", "Cache": "Cache", "Timer": "Timer"
        }
        
        result = ocr_text
        for error, correction in corrections.items():
            result = result.replace(error, correction)
        
        return result
    
    def _analyze_image_content(self, ocr_text: str, caption_text: str, slide_context: str) -> Tuple[str, List[str]]:
        """Analyze image content to determine type and extract technical keywords."""
        all_text = f"{ocr_text} {caption_text} {slide_context}".lower()
        
        # Determine image type
        image_type = "diagram"
        if any(word in all_text for word in ["chart", "graph", "performance"]):
            image_type = "chart"
        elif any(word in all_text for word in ["block", "architecture", "system"]):
            image_type = "block_diagram"
        elif any(word in all_text for word in ["flow", "process", "pipeline"]):
            image_type = "flowchart"
        
        # Extract technical keywords using similar patterns as PDF extractor
        technical_patterns = [
            r'\b(?:CPU|GPU|TPU|DSP|RISC-V|ARM|x86)\b',
            r'\b(?:I2C|SPI|UART|USB|PCIe|CAN|Ethernet)\b',
            r'\b(?:L1|L2|L3)\s*(?:cache|Cache)\b',
            r'\b(?:Timer|WDT|RAS|Cluster|Core|Cache)\b',
            r'\b(?:Ascalon|Tensix|Alexandria|Blackhole)\b',
        ]
        
        keywords = []
        for pattern in technical_patterns:
            matches = re.findall(pattern, all_text, re.IGNORECASE)
            keywords.extend([match.strip() for match in matches])
        
        return image_type, list(set([kw.strip() for kw in keywords if len(kw.strip()) > 1]))


def save_image_metadata(extracted_images: List[ExtractedImage], metadata_path: str):
    """Save image metadata to JSON file."""
    metadata = {
        "extraction_timestamp": json.dumps(None),  # Will be set by caller
        "total_images": len(extracted_images),
        "images": [img.to_dict() for img in extracted_images]
    }
    
    os.makedirs(os.path.dirname(metadata_path), exist_ok=True)
    with open(metadata_path, 'w') as f:
        json.dump(metadata, f, indent=2)
    
    logging.info(f"Saved metadata for {len(extracted_images)} images to {metadata_path}")


def load_image_metadata(metadata_path: str) -> List[ExtractedImage]:
    """Load image metadata from JSON file."""
    if not os.path.exists(metadata_path):
        return []
    
    try:
        with open(metadata_path, 'r') as f:
            metadata = json.load(f)
        
        images = []
        for img_data in metadata.get("images", []):
            images.append(ExtractedImage(**img_data))
        
        return images
        
    except Exception as e:
        logging.warning(f"Failed to load image metadata from {metadata_path}: {e}")
        return []


if __name__ == "__main__":
    # Test the image extractor
    import sys
    logging.basicConfig(level=logging.DEBUG)
    
    if len(sys.argv) > 1:
        test_file = sys.argv[1]
        
        if test_file.endswith('.pdf'):
            extractor = PDFImageExtractor()
            images = extractor.extract_images_pymupdf(test_file)
        elif test_file.endswith('.pptx'):
            extractor = PPTXImageExtractor()
            images = extractor.extract_images(test_file)
        else:
            print("Unsupported file type. Use .pdf or .pptx")
            sys.exit(1)
        
        print(f"Extracted {len(images)} images from {test_file}")
        for img in images:
            print(f"  - {img.image_path} ({img.image_type}, {len(img.technical_keywords)} keywords)")
    else:
        print("Usage: python image_extractor.py <pdf_or_pptx_file>")
