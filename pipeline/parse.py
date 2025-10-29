#!/usr/bin/env python3
"""
Parsing step for the Synapse pipeline (Parse -> Chunk -> Embed -> Retrieve -> Generate)

This script ingests heterogeneous files from a source directory (default: ../Data),
normalizes them into a unified JSONL stream of "parsed chunks" with rich metadata,
and writes them to an output file. Each line is a self-contained record suitable for
downstream chunking and embedding.

Design choices and rationale (high level):
- Keep parsing lightweight and deterministic. We extract faithful text with minimal
  transformation and attach precise locators (page, slide, heading) for citations.
- Split at natural document boundaries (page/slide/section) here so that downstream
  chunking can operate on semantically meaningful units instead of monolithic files.
- Dependencies are free and local. OCR is optional and auto-detected to avoid forcing
  system-level installs if not needed.
- NEW: Extract images/diagrams from PDFs and PPTXs as independent files with metadata

Alternate approaches and upgrades:
- unstructured.io (ENABLED by default here): Excellent element-level parsing across many
  formats with high-fidelity OCR and table extraction. Heavier dependencies. Use this when
  you care about rich structure and citations, which we do.
- Apache Tika: JVM-based, robust for enterprise doc types, but adds Java dependency.
- PDF table extraction (camelot/tabula): If you need deterministic tabular CSVs from PDFs,
  consider layering camelot/tabula. Unstructured's table HTML is often sufficient for RAG.

How to run:
  python pipeline/parse.py \
    --input "/Users/you/path/Synapse/Data" \
    --output "/Users/you/path/Synapse/artifacts/parsed.jsonl" \
    --extract-images

Outputs:
- JSONL file where each line is a dict:
  {
    "id": "sha1-of-source+locator",
    "document_id": "sha1-of-source",
    "source_path": "/abs/path/to/file.pdf",
    "source_type": "pdf|pptx|docx|md|csv|txt|image",
    "content": "extracted text",
    "metadata": { "page_number": 3, "slide_number": null, "heading_path": ["H1", "H2"], ... }
  }
- When --extract-images is used:
  - artifacts/extracted_images/ directory with image files
  - artifacts/image_metadata.json with image metadata
  - Additional image chunks in the JSONL output

From launch to query response (big picture):
- This module performs Parse only. It prepares the raw materials with citations that the
  later steps (Chunk -> Embed -> Retrieve -> Generate) will use to answer queries.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import logging
import os
import sys
import pickle
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Dict, Iterable, Iterator, List, Optional, Tuple, Callable
from datetime import datetime


# Optional/soft dependencies. Import lazily where used to keep startup fast.


# -------------------------------
# Cache management
# -------------------------------


@dataclass
class CacheEntry:
    """Cache entry storing file metadata and parsed chunks."""
    file_path: str
    file_mtime: float
    file_size: int
    chunks: List[Dict[str, object]]


def get_cache_path(output_path: str) -> str:
    """Get cache file path based on output path."""
    output_dir = os.path.dirname(output_path)
    cache_name = os.path.splitext(os.path.basename(output_path))[0] + "_cache.pkl"
    return os.path.join(output_dir, cache_name)


def load_cache(cache_path: str) -> Dict[str, CacheEntry]:
    """Load parsing cache from disk."""
    if not os.path.exists(cache_path):
        return {}
    
    try:
        with open(cache_path, 'rb') as f:
            cache = pickle.load(f)
        logging.debug(f"Loaded cache with {len(cache)} entries from {cache_path}")
        return cache
    except Exception as e:
        logging.warning(f"Failed to load cache from {cache_path}: {e}")
        return {}


def save_cache(cache: Dict[str, CacheEntry], cache_path: str) -> None:
    """Save parsing cache to disk."""
    try:
        os.makedirs(os.path.dirname(cache_path), exist_ok=True)
        with open(cache_path, 'wb') as f:
            pickle.dump(cache, f)
        logging.debug(f"Saved cache with {len(cache)} entries to {cache_path}")
    except Exception as e:
        logging.warning(f"Failed to save cache to {cache_path}: {e}")


def is_file_cached(file_path: str, cache: Dict[str, CacheEntry]) -> bool:
    """Check if file is already cached and up-to-date."""
    abs_path = os.path.abspath(file_path)
    
    if abs_path not in cache:
        return False
    
    try:
        stat = os.stat(file_path)
        entry = cache[abs_path]
        
        # Check if file modification time and size match
        return (entry.file_mtime == stat.st_mtime and 
                entry.file_size == stat.st_size)
    except (OSError, AttributeError):
        return False


def add_to_cache(file_path: str, chunks: List[ParsedChunk], cache: Dict[str, CacheEntry]) -> None:
    """Add parsed chunks to cache."""
    abs_path = os.path.abspath(file_path)
    
    try:
        stat = os.stat(file_path)
        # Convert ParsedChunk objects to dicts for serialization
        chunk_dicts = [asdict(chunk) for chunk in chunks]
        
        cache[abs_path] = CacheEntry(
            file_path=abs_path,
            file_mtime=stat.st_mtime,
            file_size=stat.st_size,
            chunks=chunk_dicts
        )
    except OSError as e:
        logging.warning(f"Failed to add {file_path} to cache: {e}")


def get_cached_chunks(file_path: str, cache: Dict[str, CacheEntry]) -> List[ParsedChunk]:
    """Get cached chunks for a file."""
    abs_path = os.path.abspath(file_path)
    entry = cache.get(abs_path)
    
    if not entry:
        return []
    
    # Convert dicts back to ParsedChunk objects
    chunks = []
    for chunk_dict in entry.chunks:
        chunk = ParsedChunk(**chunk_dict)
        chunks.append(chunk)
    
    return chunks


# -------------------------------
# Data model and utilities
# -------------------------------


@dataclass
class ParsedChunk:
    """A normalized representation of a parse unit suitable for downstream steps.

    We create one ParsedChunk per natural boundary (e.g., PDF page, PPTX slide,
    Markdown section). IDs are deterministic for reproducibility.
    """

    id: str
    document_id: str
    source_path: str
    source_type: str
    content: str
    metadata: Dict[str, object] = field(default_factory=dict)


def compute_sha1(text: str) -> str:
    hasher = hashlib.sha1()
    hasher.update(text.encode("utf-8", errors="ignore"))
    return hasher.hexdigest()


def compute_document_id(source_path: str) -> str:
    # Document ID is hash of absolute path (acts as stable, unique key within repo)
    abs_path = os.path.abspath(source_path)
    return compute_sha1(f"doc::{abs_path}")


def compute_chunk_id(document_id: str, locator: str, content: str) -> str:
    # Chunk ID includes document id + locator + content checksum to be deterministic
    content_hash = compute_sha1(content)[:16]
    return compute_sha1(f"chunk::{document_id}::{locator}::{content_hash}")


# -------------------------------
# Parsers for specific file types (basic engine)
# -------------------------------


def parse_pdf(file_path: str) -> Iterator[ParsedChunk]:
    """Parse a PDF into per-page chunks using PyMuPDF.

    - Pros: fast, robust text extraction with layout heuristics.
    - Cons: detailed table structure is not preserved; consider camelot/tabula later.
    """

    try:
        import fitz  # PyMuPDF
    except Exception as exc:  # pragma: no cover - optional dependency
        raise RuntimeError(
            "PyMuPDF (pymupdf) is required to parse PDFs. Install with 'pip install pymupdf'."
        ) from exc

    doc_id = compute_document_id(file_path)
    with fitz.open(file_path) as doc:
        page_count = doc.page_count
        for page_index in range(page_count):
            page = doc.load_page(page_index)
            text = page.get_text("text") or ""
            locator = f"page:{page_index + 1}"
            chunk_id = compute_chunk_id(doc_id, locator, text)
            yield ParsedChunk(
                id=chunk_id,
                document_id=doc_id,
                source_path=os.path.abspath(file_path),
                source_type="pdf",
                content=text,
                metadata={
                    "page_number": page_index + 1,
                    "page_count": page_count,
                    "file_name": os.path.basename(file_path),
                },
            )


def parse_pptx(file_path: str) -> Iterator[ParsedChunk]:
    """Parse a PowerPoint deck into per-slide chunks using python-pptx.

    We traverse shapes and text frames. Additionally, we extract native PPTX
    tables and serialize them to CSV for table-aware retrieval and prompting.
    """

    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover - optional dependency
        raise RuntimeError(
            "python-pptx is required to parse PPTX. Install with 'pip install python-pptx'."
        ) from exc

    # Helper to serialize a pptx table to CSV text
    def _table_to_csv(tbl) -> str:
        from io import StringIO
        import csv as _csv
        output = StringIO()
        writer = _csv.writer(output)
        try:
            for row in tbl.rows:
                cells = []
                for cell in row.cells:
                    # cell.text sometimes includes hard line breaks; normalize
                    txt = (cell.text or "").replace("\r", " ").replace("\n", " ").strip()
                    cells.append(txt)
                writer.writerow(cells)
            return output.getvalue().strip()
        finally:
            output.close()

    doc_id = compute_document_id(file_path)
    pres = Presentation(file_path)
    slide_count = len(pres.slides)

    def slide_text(slide) -> str:
        texts: List[str] = []
        for shape in slide.shapes:
            # Extract free text shapes
            if hasattr(shape, "text"):
                txt = (shape.text or "").strip()
                if txt:
                    texts.append(txt)
            elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
                frame_texts: List[str] = []
                for paragraph in shape.text_frame.paragraphs:
                    runs = [run.text for run in paragraph.runs]
                    para = "".join(runs).strip()
                    if para:
                        frame_texts.append(para)
                if frame_texts:
                    texts.append("\n".join(frame_texts))
        return "\n\n".join(texts)

    for idx, slide in enumerate(pres.slides):
        # Emit one chunk with slide narrative text
        text = slide_text(slide)
        locator = f"slide:{idx + 1}"
        chunk_id = compute_chunk_id(doc_id, locator, text)
        yield ParsedChunk(
            id=chunk_id,
            document_id=doc_id,
            source_path=os.path.abspath(file_path),
            source_type="pptx",
            content=text,
            metadata={
                "slide_number": idx + 1,
                "slide_count": slide_count,
                "file_name": os.path.basename(file_path),
            },
        )

        # Emit separate chunks for each native PPTX table on this slide (as CSV)
        table_index = 0
        for shp in slide.shapes:
            try:
                if getattr(shp, "has_table", False) and shp.table is not None:
                    table_index += 1
                    csv_text = _table_to_csv(shp.table)
                    if not csv_text.strip():
                        continue
                    t_locator = f"slide:{idx + 1}/table:{table_index}"
                    t_chunk_id = compute_chunk_id(doc_id, t_locator, csv_text)
                    yield ParsedChunk(
                        id=t_chunk_id,
                        document_id=doc_id,
                        source_path=os.path.abspath(file_path),
                        source_type="pptx",
                        content=csv_text,
                        metadata={
                            "slide_number": idx + 1,
                            "slide_count": slide_count,
                            "file_name": os.path.basename(file_path),
                            "table_index": table_index,
                            "content_format": "csv",
                            "extractor": "pptx",
                        },
                    )
            except Exception:
                continue


def parse_docx(file_path: str) -> Iterator[ParsedChunk]:
    """Parse DOCX into a single chunk (paragraphs joined)."""

    try:
        import docx  # python-docx
    except Exception as exc:  # pragma: no cover - optional dependency
        raise RuntimeError(
            "python-docx is required to parse DOCX. Install with 'pip install python-docx'."
        ) from exc

    doc_id = compute_document_id(file_path)
    d = docx.Document(file_path)
    paragraphs = [p.text.strip() for p in d.paragraphs if p.text and p.text.strip()]
    text = "\n\n".join(paragraphs)
    locator = "doc:all"
    chunk_id = compute_chunk_id(doc_id, locator, text)
    yield ParsedChunk(
        id=chunk_id,
        document_id=doc_id,
        source_path=os.path.abspath(file_path),
        source_type="docx",
        content=text,
        metadata={
            "file_name": os.path.basename(file_path),
        },
    )


def parse_markdown(file_path: str) -> Iterator[ParsedChunk]:
    """Parse Markdown splitting on ATX headings (#, ##, ###) to preserve sections.

    We keep heading hierarchy as a list in metadata.heading_path for later citation.
    """

    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()

    lines = text.splitlines()
    sections: List[Dict[str, object]] = []
    current_heading_path: List[str] = []
    current_lines: List[str] = []

    def flush_section():
        if not current_lines:
            return
        sections.append(
            {
                "heading_path": list(current_heading_path),
                "content": "\n".join(current_lines).strip(),
            }
        )

    for line in lines:
        stripped = line.strip()
        if stripped.startswith("#"):
            # New heading encountered
            level = len(stripped) - len(stripped.lstrip("#"))
            title = stripped[level:].strip()
            flush_section()
            # Adjust heading path depth
            if level <= len(current_heading_path):
                current_heading_path = current_heading_path[: level - 1]
            current_heading_path = current_heading_path + [title]
            current_lines = [f"{stripped}"]
        else:
            current_lines.append(line)

    flush_section()

    doc_id = compute_document_id(file_path)
    for idx, section in enumerate(sections):
        content = section["content"] or ""
        heading_path = section["heading_path"] or []
        locator = f"md-section:{idx + 1}"
        chunk_id = compute_chunk_id(doc_id, locator, content)
        yield ParsedChunk(
            id=chunk_id,
            document_id=doc_id,
            source_path=os.path.abspath(file_path),
            source_type="md",
            content=content,
            metadata={
                "heading_path": heading_path,
                "section_index": idx + 1,
                "file_name": os.path.basename(file_path),
            },
        )


def parse_csv(file_path: str) -> Iterator[ParsedChunk]:
    """Parse CSV by reading all rows and serializing back to CSV text.

    Justification: preserves the exact data for downstream table-aware chunkers.
    Alternative: convert to Markdown table for readability (lossy for large CSVs).
    """

    try:
        import pandas as pd
    except Exception as exc:  # pragma: no cover - optional dependency
        raise RuntimeError(
            "pandas is required to parse CSV. Install with 'pip install pandas'."
        ) from exc

    # Robust delimiter inference: try default, then python engine auto-sep, then whitespace
    def _read_csv_robust(path: str) -> "pd.DataFrame":
        try:
            return pd.read_csv(path, sep=None, engine="python")
        except Exception:
            return pd.read_csv(path, header=None, engine="python")

    df = _read_csv_robust(file_path)
    # Re-serialize to CSV text without the index
    text = df.to_csv(index=False)
    doc_id = compute_document_id(file_path)
    locator = "csv:all"
    chunk_id = compute_chunk_id(doc_id, locator, text)
    yield ParsedChunk(
        id=chunk_id,
        document_id=doc_id,
        source_path=os.path.abspath(file_path),
        source_type="csv",
        content=text,
        metadata={
            "num_rows": int(df.shape[0]),
            "num_cols": int(df.shape[1]),
            "columns": list(map(str, df.columns.tolist())),
            "file_name": os.path.basename(file_path),
            "content_format": "csv",
        },
    )


def parse_text(file_path: str) -> Iterator[ParsedChunk]:
    """Parse plain text into a single chunk."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    doc_id = compute_document_id(file_path)
    locator = "text:all"
    chunk_id = compute_chunk_id(doc_id, locator, text)
    yield ParsedChunk(
        id=chunk_id,
        document_id=doc_id,
        source_path=os.path.abspath(file_path),
        source_type="txt",
        content=text,
        metadata={"file_name": os.path.basename(file_path)},
    )


def correct_ocr_common_errors(ocr_text: str) -> str:
    """Apply common OCR error corrections for technical diagrams."""
    if not ocr_text:
        return ocr_text
    
    corrections = {
        # Common OCR errors in technical diagrams
        "12C": "I2C",
        "I3C": "I2C", 
        "SPI ": "SPI ",
        "UARTx": "UART x",
        "GPIOx": "GPIO x",
        "CANFD": "CAN-FD",
        "CANFDx": "CAN-FD x",
        "SPIx": "SPI x",
        "I2Cx": "I2C x",
        " x ": " x",  # normalize spacing around x
        "×": "x",      # replace multiplication symbol with x
        "Timer": "Timer",
        "WDT": "WDT",
    }
    
    result = ocr_text
    for error, correction in corrections.items():
        result = result.replace(error, correction)
    
    # Add context clues for better interpretation
    if any(term in result.upper() for term in ["TIMER", "WDT", "CANFD", "CAN-FD", "SPI", "I2C", "UART", "GPIO"]):
        result += "\n[Technical diagram showing system peripherals and interfaces]"
    
    return result


def parse_image_basic(
    file_path: str,
    captioner: Optional[Callable[["object"], str]] = None,
) -> Iterator[ParsedChunk]:
    """Parse image via OCR and optional captioning.

    - OCR: pytesseract if available
    - Caption: transformers image-to-text model if provided via `captioner`
    """

    try:
        import pytesseract  # type: ignore
        from PIL import Image  # type: ignore
    except Exception:  # pragma: no cover - optional dependency
        pytesseract = None  # type: ignore
        Image = None  # type: ignore

    ocr_text = ""
    caption_text = ""
    image = None
    try:
        if Image is not None:
            image = Image.open(file_path)
    except Exception as exc:
        logging.warning("Failed to open image %s: %s", file_path, exc)

    if image is not None and pytesseract is not None:
        try:
            custom_config = r'--oem 3 --psm 6 -c tessedit_char_whitelist=0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz ×x-/()[]{}:.,;'
            ocr_text = pytesseract.image_to_string(image, config=custom_config) or ""
            ocr_text = correct_ocr_common_errors(ocr_text)
        except Exception:
            ocr_text = ""

    if image is not None and captioner is not None:
        try:
            caption_text = captioner(image) or ""
        except Exception as exc:
            logging.warning("Captioning failed for %s: %s", file_path, exc)
            caption_text = ""

    if not (ocr_text.strip() or caption_text.strip()):
        return iter(())

    combined = "".join(
        part
        for part in [
            f"Caption: {caption_text.strip()}\n\n" if caption_text.strip() else "",
            f"OCR: {ocr_text.strip()}" if ocr_text.strip() else "",
        ]
    )

    doc_id = compute_document_id(file_path)
    locator = "image:caption+ocr" if caption_text.strip() else "image:ocr"
    chunk_id = compute_chunk_id(doc_id, locator, combined)
    yield ParsedChunk(
        id=chunk_id,
        document_id=doc_id,
        source_path=os.path.abspath(file_path),
        source_type="image",
        content=combined,
        metadata={
            "file_name": os.path.basename(file_path),
            "ocr_engine": "tesseract" if ocr_text.strip() else None,
            "has_caption": bool(caption_text.strip()),
        },
    )


# -------------------------------
# Unstructured.io engine (high‑fidelity with OCR & tables)
# -------------------------------


@dataclass
class UnstructuredOptions:
    pdf_strategy: str = "hi_res"  # 'hi_res' (OCR/layout), 'fast' (no OCR)
    ocr_languages: str = "eng"  # space-separated ISO 639-2 codes, e.g., "eng deu jpn"
    infer_table_structure: bool = True
    include_page_breaks: bool = True
    pdf_tables_as_csv: bool = True  # Convert unstructured Table elements to CSV content
    pdf_table_extractor: str = "pdfplumber"  # 'unstructured' or 'pdfplumber'


def _import_unstructured():
    try:  # Imports are localized so basic engine users are unaffected
        from unstructured.documents.elements import Table
        from unstructured.partition.pdf import partition_pdf
        from unstructured.partition.docx import partition_docx
        from unstructured.partition.pptx import partition_pptx
        from unstructured.partition.text import partition_text
        from unstructured.partition.md import partition_md
        from unstructured.partition.csv import partition_csv
        from unstructured.partition.image import partition_image
    except Exception as exc:  # pragma: no cover - optional dependency
        raise RuntimeError(
            "The 'unstructured' package is required for the unstructured engine. Install with 'pip install \"unstructured[all-docs]\" unstructured-inference'."
        ) from exc

    return {
        "Table": Table,
        "partition_pdf": partition_pdf,
        "partition_docx": partition_docx,
        "partition_pptx": partition_pptx,
        "partition_text": partition_text,
        "partition_md": partition_md,
        "partition_csv": partition_csv,
        "partition_image": partition_image,
    }


def _element_to_chunk(
    element,  # unstructured Element
    file_path: str,
    document_id: str,
    table_cls,
    opts: UnstructuredOptions,
) -> Optional[ParsedChunk]:
    text: str = getattr(element, "text", "") or ""
    # Skip empty whitespace-only elements
    if not text.strip() and not (hasattr(element, "text_as_html") and getattr(element, "text_as_html") is not None):
        return None

    meta = getattr(element, "metadata", None)
    page_number = getattr(meta, "page_number", None) if meta else None
    category = getattr(element, "category", None)

    metadata: Dict[str, object] = {
        "file_name": os.path.basename(file_path),
        "page_number": int(page_number) if page_number is not None else None,
        "element_type": str(category) if category else None,
    }

    # Coordinates/bounds if available
    try:
        # Different unstructured versions expose coords differently; be defensive
        if meta and getattr(meta, "coordinates", None) and getattr(meta.coordinates, "points", None):
            metadata["coordinates"] = getattr(meta.coordinates, "points")
    except Exception:
        pass

    # Handle tables with structured HTML or CSV for higher fidelity
    if table_cls is not None and isinstance(element, table_cls):
        html = getattr(element, "text_as_html", None)
        if html:
            metadata["table_html"] = html
            if opts.pdf_tables_as_csv:
                csv_text = _html_table_to_csv(html)
                if csv_text.strip():
                    text = csv_text
                    metadata["content_format"] = "csv"
                    metadata["extractor"] = "unstructured"
                    metadata["source_table_format"] = "html"

    # Locator uses page + category + index hash via content
    locator_parts: List[str] = []
    if page_number is not None:
        locator_parts.append(f"page:{int(page_number)}")
    if category:
        locator_parts.append(f"elt:{category}")
    locator = "/".join(locator_parts) if locator_parts else "elt"
    chunk_id = compute_chunk_id(document_id, locator, text if text else metadata.get("table_html", "") or "")

    return ParsedChunk(
        id=chunk_id,
        document_id=document_id,
        source_path=os.path.abspath(file_path),
        source_type=os.path.splitext(file_path)[1].lstrip(".").lower(),
        content=text,
        metadata=metadata,
    )


def parse_file_unstructured(
    file_path: str,
    opts: UnstructuredOptions,
    image_captioner: Optional[Callable[["object"], str]] = None,
) -> List[ParsedChunk]:
    u = _import_unstructured()
    ext = os.path.splitext(file_path)[1].lower()
    document_id = compute_document_id(file_path)

    elements = []
    try:
        if ext == ".pdf":
            elements = u["partition_pdf"](
                filename=file_path,
                strategy=opts.pdf_strategy,
                ocr_languages=opts.ocr_languages,
                infer_table_structure=opts.infer_table_structure,
                include_page_breaks=opts.include_page_breaks,
            )
        elif ext == ".pptx":
            elements = u["partition_pptx"](filename=file_path, infer_table_structure=opts.infer_table_structure)
        elif ext == ".docx":
            elements = u["partition_docx"](filename=file_path)
        elif ext == ".md":
            elements = u["partition_md"](filename=file_path)
        elif ext == ".csv":
            elements = u["partition_csv"](filename=file_path)
        elif ext in {".txt", ".log"}:
            elements = u["partition_text"](filename=file_path)
        elif ext in {".png", ".jpg", ".jpeg", ".tiff", ".bmp"}:
            # Prefer a single combined chunk with caption + OCR for images
            try:
                elements = u["partition_image"](filename=file_path, ocr_languages=opts.ocr_languages)
            except Exception:
                elements = []
            ocr_text = "\n".join((getattr(el, "text", "") or "") for el in elements).strip()
            
            # Apply OCR error corrections to unstructured results
            ocr_text = correct_ocr_common_errors(ocr_text)
            
            caption_text = ""
            try:
                from PIL import Image  # type: ignore
                img = Image.open(file_path)
            except Exception:
                img = None
            if img is not None and image_captioner is not None:
                try:
                    caption_text = image_captioner(img) or ""
                except Exception as exc:
                    logging.warning("Captioning failed for %s: %s", file_path, exc)
            combined = "".join(
                part
                for part in [
                    f"Caption: {caption_text}\n\n" if caption_text else "",
                    f"OCR: {ocr_text}" if ocr_text else "",
                ]
            )
            if combined.strip():
                locator = "image:caption+ocr" if caption_text else "image:ocr"
                chunk_id = compute_chunk_id(document_id, locator, combined)
                return [
                    ParsedChunk(
                        id=chunk_id,
                        document_id=document_id,
                        source_path=os.path.abspath(file_path),
                        source_type="image",
                        content=combined,
                        metadata={
                            "file_name": os.path.basename(file_path),
                            "has_caption": bool(caption_text),
                            "ocr_engine": "tesseract" if ocr_text else None,
                        },
                    )
                ]
            return []
        else:
            logging.info("Unstructured: skipping unsupported file type: %s", file_path)
            return []
    except Exception as exc:
        logging.warning("Unstructured failed for %s: %s", file_path, exc)
        return []

    chunks: List[ParsedChunk] = []
    table_cls = u["Table"]

    # If we choose pdfplumber as the table extractor, we skip unstructured Table elements to avoid duplication
    for el in elements:
        if ext == ".pdf" and opts.pdf_table_extractor == "pdfplumber" and isinstance(el, table_cls):
            continue
        chunk = _element_to_chunk(el, file_path, document_id, table_cls, opts)
        if chunk is not None:
            chunks.append(chunk)

    # Add tables from pdfplumber if requested
    if ext == ".pdf" and opts.pdf_table_extractor == "pdfplumber":
        plumber_chunks = extract_pdf_tables_with_pdfplumber(file_path, document_id)
        chunks.extend(plumber_chunks)

    return chunks


def _html_table_to_csv(html: str) -> str:
    """Convert HTML table(s) to CSV. Returns a CSV string.

    We use pandas.read_html (lxml backend). If multiple tables are found within the HTML,
    we concatenate them separated by a blank line to retain all content.
    """
    try:
        import pandas as pd
    except Exception:
        return ""

    try:
        dfs = pd.read_html(html)  # type: ignore[arg-type]
        if not dfs:
            return ""
        csv_parts: List[str] = []
        for df in dfs:
            csv_parts.append(df.to_csv(index=False))
        return "\n".join(csv_parts)
    except Exception:
        return ""


def extract_pdf_tables_with_pdfplumber(file_path: str, document_id: str) -> List[ParsedChunk]:
    """Extract tables from a PDF using pdfplumber and return them as CSV chunks.

    Strategy: Try line-based (lattice-like) detection, fallback to text-based (stream) if none.
    """
    try:
        import pdfplumber  # type: ignore
        import pandas as pd  # type: ignore
    except Exception as exc:  # pragma: no cover - optional dependency
        logging.warning("pdfplumber not available for table extraction: %s", exc)
        return []

    def rows_to_dataframe(rows: List[List[Optional[str]]]) -> "pd.DataFrame":
        # Heuristic: use first non-empty row as header if plausible
        import pandas as pd  # local import
        if not rows:
            return pd.DataFrame()
        header_idx = 0
        while header_idx < len(rows) and all((c is None or str(c).strip() == "") for c in rows[header_idx]):
            header_idx += 1
        data_rows = rows[header_idx + 1 :] if header_idx < len(rows) - 1 else []
        header = rows[header_idx] if header_idx < len(rows) else []
        if header and len(set(map(str, header))) == len(header):
            df = pd.DataFrame(data_rows, columns=[str(c) if c is not None else "" for c in header])
        else:
            df = pd.DataFrame(rows)
        return df

    chunks: List[ParsedChunk] = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for page_idx, page in enumerate(pdf.pages):
                tables: List[List[List[Optional[str]]]] = []
                tables = page.extract_tables() or []

                for t_idx, rows in enumerate(tables):
                    try:
                        df = rows_to_dataframe(rows)
                        if df.empty:
                            continue
                        csv_text = df.to_csv(index=False)
                        locator = f"page:{page_idx + 1}/table:{t_idx + 1}"
                        chunk_id = compute_chunk_id(document_id, locator, csv_text)
                        chunks.append(
                            ParsedChunk(
                                id=chunk_id,
                                document_id=document_id,
                                source_path=os.path.abspath(file_path),
                                source_type="pdf",
                                content=csv_text,
                                metadata={
                                    "file_name": os.path.basename(file_path),
                                    "page_number": page_idx + 1,
                                    "table_index": t_idx + 1,
                                    "extractor": "pdfplumber",
                                },
                            )
                        )
                    except Exception:
                        continue
    except Exception as exc:
        logging.warning("Failed pdfplumber extraction for %s: %s", file_path, exc)
        return chunks

    return chunks


# -------------------------------
# Driver that walks the directory
# -------------------------------


BASIC_SUPPORTED_EXTENSIONS = {
    ".pdf": parse_pdf,
    ".pptx": parse_pptx,
    ".docx": parse_docx,
    ".md": parse_markdown,
    ".csv": parse_csv,
    ".txt": parse_text,
    # Images handled specially in parse_file_basic to allow captioning
}


def iter_files(root_dir: str) -> Iterator[str]:
    for dirpath, dirnames, filenames in os.walk(root_dir):
        # Skip hidden directories like .git, .DS_Store
        dirnames[:] = [d for d in dirnames if not d.startswith(".")]
        for filename in filenames:
            if filename.startswith("."):
                continue
            yield os.path.join(dirpath, filename)


def get_folder_structure(root_dir: str) -> Dict[str, List[str]]:
    """Get folder-based organization for embeddings.
    
    Returns:
        Dict mapping folder_key -> list of file paths
        
    Rules:
    - Regular folders: create one embedding per folder
    - Folders starting with '#': dig deeper, create embedding per subfolder
    """
    folder_groups: Dict[str, List[str]] = {}
    root_path = Path(root_dir)
    
    for dirpath, dirnames, filenames in os.walk(root_dir):
        # Skip hidden directories
        dirnames[:] = [d for d in dirnames if not d.startswith(".")]
        
        current_path = Path(dirpath)
        relative_path = current_path.relative_to(root_path)
        
        # Skip if no files in this directory
        valid_files = [f for f in filenames if not f.startswith(".")]
        if not valid_files:
            continue
            
        # Determine folder key based on structure
        parts = relative_path.parts
        
        if len(parts) == 0:
            # Files in root directory
            folder_key = "root"
        elif len(parts) == 1:
            # First level directories
            if parts[0].startswith('#'):
                # Skip - we'll handle subfolders of # directories
                continue
            else:
                folder_key = parts[0]
        else:
            # Deeper directories
            if parts[0].startswith('#'):
                # For # directories, use the subfolder as the key
                folder_key = f"{parts[0]}/{parts[1]}"
            else:
                # For regular deep directories, use the top-level folder
                folder_key = parts[0]
        
        # Add files to the appropriate group
        if folder_key not in folder_groups:
            folder_groups[folder_key] = []
            
        for filename in valid_files:
            full_path = os.path.join(dirpath, filename)
            folder_groups[folder_key].append(full_path)
    
    return folder_groups


def parse_file_basic(file_path: str, image_captioner: Optional[Callable[["object"], str]] = None) -> List[ParsedChunk]:
    ext = os.path.splitext(file_path)[1].lower()
    if ext in {".png", ".jpg", ".jpeg", ".tiff", ".bmp"}:
        try:
            return list(parse_image_basic(file_path, captioner=image_captioner))
        except Exception as exc:
            logging.warning("Failed to parse image %s: %s", file_path, exc)
            return []
    parser = BASIC_SUPPORTED_EXTENSIONS.get(ext)
    if parser is None:
        logging.info("Skipping unsupported file type: %s", file_path)
        return []
    try:
        return list(parser(file_path))
    except Exception as exc:
        logging.warning("Failed to parse %s: %s", file_path, exc)
        return []


def write_jsonl(chunks: Iterable[ParsedChunk], out_path: str) -> None:
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    count = 0
    with open(out_path, "w", encoding="utf-8") as f:
        for chunk in chunks:
            record = asdict(chunk)
            f.write(json.dumps(record, ensure_ascii=False) + "\n")
            count += 1
    logging.info("Wrote %d parsed chunks to %s", count, out_path)


def extract_images_from_file(file_path: str, pdf_extractor, pptx_extractor, parsed_chunks: List[ParsedChunk]):
    """Extract images from a file based on its type."""
    try:
        from pipeline.image_extractor import ExtractedImage
    except ImportError:
        # If running from the Synapse directory, try relative import
        pipeline_path = os.path.join(os.path.dirname(__file__))
        if pipeline_path not in sys.path:
            sys.path.insert(0, pipeline_path)
        from image_extractor import ExtractedImage
    
    file_ext = os.path.splitext(file_path)[1].lower()
    extracted_images = []
    
    try:
        if file_ext == ".pdf" and pdf_extractor:
            # Build page context from parsed chunks for better image metadata
            page_context = {}
            for chunk in parsed_chunks:
                if chunk.source_type == "pdf":
                    page_num = chunk.metadata.get("page_number")
                    if page_num:
                        page_text = page_context.get(page_num, "")
                        page_context[page_num] = page_text + " " + chunk.content
            
            extracted_images = pdf_extractor.extract_images_pymupdf(file_path, page_context)
            
        elif file_ext == ".pptx" and pptx_extractor:
            extracted_images = pptx_extractor.extract_images(file_path)
            
    except Exception as e:
        logging.warning("Failed to extract images from %s: %s", file_path, e)
    
    return extracted_images


def convert_images_to_chunks(extracted_images) -> List[ParsedChunk]:
    """Convert ExtractedImage objects to ParsedChunk objects for the pipeline."""
    chunks = []
    
    for img in extracted_images:
        try:
            # Create a comprehensive text description combining all image information
            content_parts = []
            
            # Add image type and basic info
            content_parts.append(f"Image Type: {img.image_type}")
            content_parts.append(f"Dimensions: {img.width}x{img.height} pixels")
            
            # Add AI-generated caption if available
            if img.caption_text:
                content_parts.append(f"Caption: {img.caption_text}")
            
            # Add OCR text if available
            if img.ocr_text:
                content_parts.append(f"OCR Text: {img.ocr_text}")
            
            # Add technical keywords
            if img.technical_keywords:
                content_parts.append(f"Technical Keywords: {', '.join(img.technical_keywords)}")
            
            # Add document context (truncated)
            if img.document_context:
                context_preview = img.document_context[:500] + "..." if len(img.document_context) > 500 else img.document_context
                content_parts.append(f"Document Context: {context_preview}")
            
            content = "\n\n".join(content_parts)
            
            # Create locator for the image
            locator = f"image:{img.page_number or 'unknown'}"
            
            # Compute chunk ID
            chunk_id = compute_chunk_id(img.source_document, locator, content)
            
            # Create metadata
            metadata = {
                "page_number": img.page_number,
                "file_name": os.path.basename(img.source_document),
                "image_path": img.image_path,
                "image_id": img.image_id,
                "image_type": img.image_type,
                "image_format": img.image_format,
                "width": img.width,
                "height": img.height,
                "extraction_method": img.extraction_method,
                "technical_keywords": img.technical_keywords,
                "has_ocr": bool(img.ocr_text),
                "has_caption": bool(img.caption_text),
                "content_format": "extracted_image"
            }
            
            chunk = ParsedChunk(
                id=chunk_id,
                document_id=compute_document_id(img.source_document),
                source_path=img.source_document,
                source_type="image",  # Mark as image type for special handling
                content=content,
                metadata=metadata
            )
            
            chunks.append(chunk)
            
        except Exception as e:
            logging.warning("Failed to convert extracted image to chunk: %s", e)
    
    return chunks


def main(argv: Optional[List[str]] = None) -> int:
    parser = argparse.ArgumentParser(description="Parse documents into normalized JSONL chunks")
    default_input = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "Data"))
    default_output_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "artifacts"))
    parser.add_argument("--input", type=str, default=default_input, help="Input directory (root of documents)")
    parser.add_argument(
        "--output",
        type=str,
        default=os.path.join(default_output_dir, "parsed.jsonl"),
        help="Output JSONL path (will be suffixed with folder names for folder-based parsing)",
    )
    parser.add_argument(
        "--folder-based",
        action="store_true",
        help="Create separate parsed files for each folder (enables specialized knowledge bases)",
    )
    parser.add_argument(
        "--engine",
        type=str,
        choices=["unstructured", "basic"],
        default="basic",
        help="Parsing engine to use. 'basic' uses lightweight parsers (fast); 'unstructured' enables high-fidelity OCR & tables (slow).",
    )
    parser.add_argument("--pdf-strategy", type=str, default="hi_res", choices=["hi_res", "fast"], help="Unstructured PDF strategy: 'hi_res' uses OCR/layout for scanned/complex docs; 'fast' is text-first")
    parser.add_argument("--ocr-languages", type=str, default="eng", help="OCR languages (space-separated, e.g., 'eng deu jpn') for unstructured")
    parser.add_argument(
        "--no-infer-tables",
        action="store_true",
        help="Disable table structure inference in unstructured (keeps text only)",
    )
    parser.add_argument(
        "--pdf-table-extractor",
        type=str,
        choices=["unstructured", "pdfplumber"],
        default="pdfplumber",
        help="How to extract PDF tables as CSV. 'pdfplumber' is usually best for engineering docs.",
    )
    parser.add_argument(
        "--no-pdf-tables-as-csv",
        action="store_true",
        help="Do not convert unstructured Table elements to CSV (keep HTML/text)",
    )
    parser.add_argument("--verbose", action="store_true", help="Verbose logging")
    parser.add_argument(
        "--image-captioning",
        type=str,
        choices=["off", "blip"],
        default="blip",
        help="Enable image captioning and choose backend (default: blip)",
    )
    parser.add_argument(
        "--image-captioner-model",
        type=str,
        default="Salesforce/blip-image-captioning-base",
        help="Model to use for image captioning",
    )
    parser.add_argument(
        "--image-caption-max-new-tokens",
        type=int,
        default=64,
        help="Max new tokens for caption generation",
    )
    parser.add_argument(
        "--extract-images",
        action="store_true",
        help="Extract images/diagrams from PDFs and PPTXs as independent files",
    )
    parser.add_argument(
        "--images-output-dir",
        type=str,
        default=os.path.join(default_output_dir, "extracted_images"),
        help="Directory to save extracted images",
    )
    parser.add_argument(
        "--enable-image-captioning", 
        action="store_true", 
        help="Enable AI captioning for extracted images (slower, may hit rate limits)"
    )
    parser.add_argument(
        "--disable-image-deduplication",
        action="store_true",
        help="Disable image deduplication (will save all images even if identical)"
    )
    parser.add_argument(
        "--image-similarity-threshold",
        type=int,
        default=10,
        help="Similarity threshold for perceptual image hashing (0-64, lower = more strict). Default: 10"
    )
    parser.add_argument(
        "--fast-image-extraction",
        action="store_true",
        help="Enable fast image extraction mode (disables captioning, uses simpler OCR, faster processing)"
    )
    parser.add_argument(
        "--no-cache",
        action="store_true",
        help="Disable caching - reparse all files even if they haven't changed",
    )
    parser.add_argument(
        "--cache-path",
        type=str,
        help="Custom cache file path (default: auto-generated based on output path)",
    )
    args = parser.parse_args(argv)

    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.INFO,
        format="%(asctime)s %(levelname)s %(message)s",
    )

    input_dir = os.path.abspath(args.input)
    output_path = os.path.abspath(args.output)

    if not os.path.isdir(input_dir):
        logging.error("Input directory does not exist: %s", input_dir)
        return 2

    logging.info("Parsing input directory: %s", input_dir)

    # Build optional image captioner
    image_captioner = None
    if args.image_captioning != "off":
        try:
            from transformers import pipeline  # type: ignore

            pipe = pipeline("image-to-text", model=args.image_captioner_model)

            def _caption(image) -> str:
                out = pipe(image, max_new_tokens=int(args.image_caption_max_new_tokens))
                if isinstance(out, list) and out and isinstance(out[0], dict) and "generated_text" in out[0]:
                    return str(out[0]["generated_text"]).strip()
                return ""

            image_captioner = _caption
        except Exception as exc:
            logging.warning("Image captioner unavailable: %s", exc)
    
    use_unstructured = args.engine == "unstructured"
    uopts = UnstructuredOptions(
        pdf_strategy=args.pdf_strategy,
        ocr_languages=args.ocr_languages,
        infer_table_structure=not args.no_infer_tables,
        include_page_breaks=True,
        pdf_tables_as_csv=not args.no_pdf_tables_as_csv,
    )

    if args.folder_based:
        # Folder-based parsing: create separate files for each folder
        folder_groups = get_folder_structure(input_dir)
        logging.info("Folder-based parsing enabled. Found %d folder groups:", len(folder_groups))
        for folder_key in sorted(folder_groups.keys()):
            file_count = len(folder_groups[folder_key])
            logging.info("  - %s: %d files", folder_key, file_count)
        
        return process_folder_based_parsing(
            folder_groups, args, input_dir, output_path, 
            use_unstructured, uopts, image_captioner
        )
    else:
        # Original single-file parsing
        return process_single_file_parsing(
            input_dir, args, output_path, use_unstructured, uopts, image_captioner
        )


def process_single_file_parsing(
    input_dir: str, args, output_path: str, use_unstructured: bool, 
    uopts, image_captioner
) -> int:
    """Original single-file parsing logic."""
    # Initialize cache
    cache_path = args.cache_path or get_cache_path(output_path)
    cache: Dict[str, CacheEntry] = {} if args.no_cache else load_cache(cache_path)
    
    all_chunks: List[ParsedChunk] = []
    cached_count = 0
    parsed_count = 0
    
    # Initialize progress tracking
    all_files = list(iter_files(input_dir))
    progress_tracker = None
    if not args.verbose:  # Only show progress bar in non-verbose mode
        try:
            from .progress import ProgressTracker
            progress_tracker = ProgressTracker(len(all_files), "Parsing")
        except ImportError:
            pass
    
    # Initialize image extraction if requested
    all_extracted_images = []
    pdf_extractor = None
    pptx_extractor = None
    image_dedup_cache = None
    if args.extract_images:
        try:
            # Try different import paths for the image extractor
            try:
                from pipeline.image_extractor import PDFImageExtractor, PPTXImageExtractor, save_image_metadata, ImageDeduplicationCache
            except ImportError:
                # If running from the Synapse directory, try relative import
                pipeline_path = os.path.join(os.path.dirname(__file__))
                if pipeline_path not in sys.path:
                    sys.path.insert(0, pipeline_path)
                from image_extractor import PDFImageExtractor, PPTXImageExtractor, save_image_metadata, ImageDeduplicationCache
            
            # Create a shared deduplication cache for all files (unless disabled)
            image_dedup_cache = None if args.disable_image_deduplication else ImageDeduplicationCache(similarity_threshold=args.image_similarity_threshold)
            
            pdf_extractor = PDFImageExtractor(output_dir=args.images_output_dir, enable_captioning=args.enable_image_captioning, dedup_cache=image_dedup_cache, fast_mode=args.fast_image_extraction)
            pptx_extractor = PPTXImageExtractor(output_dir=args.images_output_dir, enable_captioning=args.enable_image_captioning, dedup_cache=image_dedup_cache, fast_mode=args.fast_image_extraction)
            
            dedup_status = "with deduplication" if not args.disable_image_deduplication else "without deduplication"
            logging.info(f"Image extraction enabled {dedup_status} - will extract images from PDFs and PPTXs")
        except ImportError as e:
            logging.warning("Image extraction disabled due to missing dependencies: %s", e)
            args.extract_images = False

    for file_idx, file_path in enumerate(all_files):
        # Update progress
        if progress_tracker:
            progress_tracker.update(file_path, file_idx)
        
        # Check cache first
        if not args.no_cache and is_file_cached(file_path, cache):
            chunks = get_cached_chunks(file_path, cache)
            cached_count += 1
            logging.debug("Using cached %d chunks from %s", len(chunks), file_path)
        else:
            # Parse the file
            chunks: List[ParsedChunk]
            if use_unstructured:
                chunks = parse_file_unstructured(file_path, uopts, image_captioner=image_captioner)
            else:
                chunks = parse_file_basic(file_path, image_captioner=image_captioner)
            
            parsed_count += 1
            if chunks:
                logging.debug("Parsed %d chunks from %s", len(chunks), file_path)
                # Add to cache if caching is enabled
                if not args.no_cache:
                    add_to_cache(file_path, chunks, cache)
        
        # Extract images if requested
        if args.extract_images:
            extracted_images = extract_images_from_file(
                file_path, pdf_extractor, pptx_extractor, chunks
            )
            if extracted_images:
                all_extracted_images.extend(extracted_images)
                # Convert extracted images to ParsedChunk objects and add to chunks
                image_chunks = convert_images_to_chunks(extracted_images)
                chunks.extend(image_chunks)
                logging.debug("Extracted %d images from %s", len(extracted_images), file_path)
        
        all_chunks.extend(chunks)
    
    # Finish progress tracking
    if progress_tracker:
        progress_tracker.finish()

    # Sort for reproducibility (by source then locator-ish via id)
    all_chunks.sort(key=lambda c: (c.source_path, c.id))
    write_jsonl(all_chunks, output_path)
    
    # Save image metadata if images were extracted
    if args.extract_images and all_extracted_images:
        metadata_path = os.path.join(os.path.dirname(output_path), "image_metadata.json")
        save_image_metadata(all_extracted_images, metadata_path)
        
        # Count unique vs duplicate images
        unique_images = len([img for img in all_extracted_images if not img.is_duplicate])
        duplicate_images = len([img for img in all_extracted_images if img.is_duplicate])
        
        logging.info("Saved metadata for %d extracted images to %s", len(all_extracted_images), metadata_path)
        logging.info("Image deduplication: %d unique images, %d duplicates avoided", unique_images, duplicate_images)
    
    # Save cache
    if not args.no_cache:
        save_cache(cache, cache_path)
    
    # Print parsing summary with big font
    print("\n" + "="*80)
    print("██████╗  █████╗ ██████╗ ███████╗██╗███╗   ██╗ ██████╗ ")
    print("██╔══██╗██╔══██╗██╔══██╗██╔════╝██║████╗  ██║██╔════╝ ")
    print("██████╔╝███████║██████╔╝███████╗██║██╔██╗ ██║██║  ███╗")
    print("██╔═══╝ ██╔══██║██╔══██╗╚════██║██║██║╚██╗██║██║   ██║")
    print("██║     ██║  ██║██║  ██║███████║██║██║ ╚████║╚██████╔╝")
    print("╚═╝     ╚═╝  ╚═╝╚═╝  ╚═╝╚══════╝╚═╝╚═╝  ╚═══╝ ╚═════╝ ")
    print("                    ███████╗██╗   ██╗███╗   ███╗███╗   ███╗ █████╗ ██████╗ ██╗   ██╗")
    print("                    ██╔════╝██║   ██║████╗ ████║████╗ ████║██╔══██╗██╔══██╗╚██╗ ██╔╝")
    print("                    ███████╗██║   ██║██╔████╔██║██╔████╔██║███████║██████╔╝ ╚████╔╝ ")
    print("                    ╚════██║██║   ██║██║╚██╔╝██║██║╚██╔╝██║██╔══██║██╔══██╗  ╚██╔╝  ")
    print("                    ███████║╚██████╔╝██║ ╚═╝ ██║██║ ╚═╝ ██║██║  ██║██║  ██║   ██║   ")
    print("                    ╚══════╝ ╚═════╝ ╚═╝     ╚═╝╚═╝     ╚═╝╚═╝  ╚═╝╚═╝  ╚═╝   ╚═╝   ")
    print("="*80)
    print(f"📁 INPUT DIRECTORY: {input_dir}")
    print(f"📄 OUTPUT FILE: {output_path}")
    print(f"🔢 TOTAL CHUNKS: {len(all_chunks):,}")
    
    # Cache statistics
    total_files = parsed_count + cached_count
    if not args.no_cache and total_files > 0:
        print(f"⚡ CACHE STATISTICS:")
        print(f"   📋 Cache file: {cache_path}")
        print(f"   🔄 Files parsed: {parsed_count}")
        print(f"   💾 Files from cache: {cached_count}")
        print(f"   📈 Cache hit rate: {cached_count/total_files*100:.1f}%")
    elif args.no_cache:
        print(f"🚫 CACHING DISABLED - All {total_files} files were parsed")
    
    # Count by file type
    type_counts = {}
    for chunk in all_chunks:
        file_type = chunk.source_type.upper()
        type_counts[file_type] = type_counts.get(file_type, 0) + 1
    
    print(f"📊 BREAKDOWN BY TYPE:")
    for file_type, count in sorted(type_counts.items()):
        print(f"   {file_type}: {count:,} chunks")
    
    # Count by unique documents
    unique_docs = len(set(chunk.document_id for chunk in all_chunks))
    print(f"📚 UNIQUE DOCUMENTS: {unique_docs}")
    
    # Image extraction summary
    if args.extract_images:
        print(f"🖼️  IMAGE EXTRACTION:")
        if all_extracted_images:
            unique_images = len([img for img in all_extracted_images if not img.is_duplicate])
            duplicate_images = len([img for img in all_extracted_images if img.is_duplicate])
            
            print(f"   📷 Total images processed: {len(all_extracted_images)}")
            print(f"   ✨ Unique images saved: {unique_images}")
            print(f"   🔗 Duplicate references: {duplicate_images}")
            print(f"   💾 Images saved to: {args.images_output_dir}")
            print(f"   📋 Metadata saved to: {os.path.join(os.path.dirname(output_path), 'image_metadata.json')}")
            
            # Break down by image type
            image_types = {}
            for img in all_extracted_images:
                img_type = img.image_type
                image_types[img_type] = image_types.get(img_type, 0) + 1
            
            print(f"   📊 Images by type:")
            for img_type, count in sorted(image_types.items()):
                print(f"      {img_type}: {count}")
                
            if duplicate_images > 0:
                space_saved = duplicate_images  # Rough estimate - each duplicate saves one file
                print(f"   💡 Deduplication saved ~{space_saved} duplicate files")
        else:
            print(f"   📷 No images extracted (may be none found or extraction failed)")
    else:
        print(f"🖼️  IMAGE EXTRACTION: Disabled (use --extract-images to enable)")
    
    # Show unstructured command for higher quality processing
    if args.engine == "basic":
        print(f"🔬 FOR HIGHER QUALITY (OCR + ADVANCED TABLES):")
        unstructured_cmd = f"python3 pipeline/parse.py --engine unstructured --input \"{input_dir}\" --output \"{output_path.replace('.jsonl', '_unstructured.jsonl')}\""
        print(f"   pip3 install \"unstructured[all-docs]\" unstructured-inference")
        print(f"   {unstructured_cmd}")
        print(f"   ⚠️  Note: Unstructured processing takes 10-50x longer but provides OCR for scanned PDFs")
    
    print("="*80)
    print("✅ PARSING COMPLETED SUCCESSFULLY!")
    print("="*80 + "\n")
    
    logging.info("Parsing completed successfully.")
    return 0


def process_folder_based_parsing(
    folder_groups: Dict[str, List[str]], args, input_dir: str, 
    base_output_path: str, use_unstructured: bool, uopts, image_captioner
) -> int:
    """Process files with folder-based organization, creating separate artifacts per folder."""
    
    # Initialize progress tracking for folder-based parsing
    total_files = sum(len(files) for files in folder_groups.values())
    progress_tracker = None
    if not args.verbose:  # Only show progress bar in non-verbose mode
        try:
            from .progress import ProgressTracker
            progress_tracker = ProgressTracker(total_files, "Parsing (folder-based)")
        except ImportError:
            pass
    
    # Initialize image extraction if requested
    pdf_extractor = None
    pptx_extractor = None
    image_dedup_cache = None
    if args.extract_images:
        try:
            # Try different import paths for the image extractor
            try:
                from pipeline.image_extractor import PDFImageExtractor, PPTXImageExtractor, save_image_metadata, ImageDeduplicationCache
            except ImportError:
                # If running from the Synapse directory, try relative import
                pipeline_path = os.path.join(os.path.dirname(__file__))
                if pipeline_path not in sys.path:
                    sys.path.insert(0, pipeline_path)
                from image_extractor import PDFImageExtractor, PPTXImageExtractor, save_image_metadata, ImageDeduplicationCache
            
            # Create a shared deduplication cache for all files across all folders (unless disabled)
            image_dedup_cache = None if args.disable_image_deduplication else ImageDeduplicationCache(similarity_threshold=args.image_similarity_threshold)
            
            pdf_extractor = PDFImageExtractor(output_dir=args.images_output_dir, enable_captioning=args.enable_image_captioning, dedup_cache=image_dedup_cache, fast_mode=args.fast_image_extraction)
            pptx_extractor = PPTXImageExtractor(output_dir=args.images_output_dir, enable_captioning=args.enable_image_captioning, dedup_cache=image_dedup_cache, fast_mode=args.fast_image_extraction)
            
            dedup_status = "with deduplication" if not args.disable_image_deduplication else "without deduplication"
            logging.info(f"Image extraction enabled {dedup_status} - will extract images from PDFs and PPTXs")
        except ImportError as e:
            logging.warning("Image extraction disabled due to missing dependencies: %s", e)
            args.extract_images = False
    
    total_chunks = 0
    total_cached = 0
    total_parsed = 0
    all_folder_summaries = []
    
    # Process each folder group
    for folder_key, file_paths in folder_groups.items():
        logging.info("Processing folder: %s (%d files)", folder_key, len(file_paths))
        
        # Create folder-specific output path
        safe_folder_name = folder_key.replace('/', '_').replace('#', 'hash_')
        folder_output_path = base_output_path.replace('.jsonl', f'_{safe_folder_name}.jsonl')
        
        # Create folder-specific cache
        cache_path = args.cache_path or get_cache_path(folder_output_path)
        cache: Dict[str, CacheEntry] = {} if args.no_cache else load_cache(cache_path)
        
        folder_chunks: List[ParsedChunk] = []
        folder_extracted_images = []
        cached_count = 0
        parsed_count = 0
        
        # Process files in this folder
        for file_path in file_paths:
            # Update progress
            if progress_tracker:
                progress_tracker.update(f"{folder_key}: {os.path.basename(file_path)}")
            
            # Check cache first
            if not args.no_cache and is_file_cached(file_path, cache):
                chunks = get_cached_chunks(file_path, cache)
                cached_count += 1
                logging.debug("Using cached %d chunks from %s", len(chunks), file_path)
            else:
                # Parse the file
                chunks: List[ParsedChunk]
                if use_unstructured:
                    chunks = parse_file_unstructured(file_path, uopts, image_captioner=image_captioner)
                else:
                    chunks = parse_file_basic(file_path, image_captioner=image_captioner)
                
                parsed_count += 1
                if chunks:
                    logging.debug("Parsed %d chunks from %s", len(chunks), file_path)
                    # Add to cache if caching is enabled
                    if not args.no_cache:
                        add_to_cache(file_path, chunks, cache)
            
            # Extract images if requested
            if args.extract_images:
                extracted_images = extract_images_from_file(
                    file_path, pdf_extractor, pptx_extractor, chunks
                )
                if extracted_images:
                    folder_extracted_images.extend(extracted_images)
                    # Convert extracted images to ParsedChunk objects and add to chunks
                    image_chunks = convert_images_to_chunks(extracted_images)
                    chunks.extend(image_chunks)
                    logging.debug("Extracted %d images from %s", len(extracted_images), file_path)
            
            folder_chunks.extend(chunks)
        
        # Sort and write folder chunks
        folder_chunks.sort(key=lambda c: (c.source_path, c.id))
        write_jsonl(folder_chunks, folder_output_path)
        
        # Save folder-specific image metadata
        if args.extract_images and folder_extracted_images:
            metadata_path = os.path.join(os.path.dirname(folder_output_path), f"image_metadata_{safe_folder_name}.json")
            save_image_metadata(folder_extracted_images, metadata_path)
            logging.info("Saved metadata for %d extracted images to %s", len(folder_extracted_images), metadata_path)
        
        # Save folder cache
        if not args.no_cache:
            save_cache(cache, cache_path)
        
        # Track folder summary
        folder_summary = {
            'folder_key': folder_key,
            'safe_name': safe_folder_name,
            'output_path': folder_output_path,
            'chunks': len(folder_chunks),
            'files_parsed': parsed_count,
            'files_cached': cached_count,
            'images_extracted': len(folder_extracted_images) if folder_extracted_images else 0
        }
        all_folder_summaries.append(folder_summary)
        
        total_chunks += len(folder_chunks)
        total_cached += cached_count
        total_parsed += parsed_count
        
        logging.info("Completed folder %s: %d chunks, %d files", folder_key, len(folder_chunks), len(file_paths))
    
    # Finish progress tracking
    if progress_tracker:
        progress_tracker.finish()
    
    # Print comprehensive summary
    print("\n" + "="*80)
    print("██████╗  █████╗ ██████╗ ███████╗██╗███╗   ██╗ ██████╗ ")
    print("██╔══██╗██╔══██╗██╔══██╗██╔════╝██║████╗  ██║██╔════╝ ")
    print("██████╔╝███████║██████╔╝███████╗██║██╔██╗ ██║██║  ███╗")
    print("██╔═══╝ ██╔══██║██╔══██╗╚════██║██║██║╚██╗██║██║   ██║")
    print("██║     ██║  ██║██║  ██║███████║██║██║ ╚████║╚██████╔╝")
    print("╚═╝     ╚═╝  ╚═╝╚═╝  ╚═╝╚══════╝╚═╝╚═╝  ╚═══╝ ╚═════╝ ")
    print("                    ███████╗██╗   ██╗███╗   ███╗███╗   ███╗ █████╗ ██████╗ ██╗   ██╗")
    print("                    ██╔════╝██║   ██║████╗ ████║████╗ ████║██╔══██╗██╔══██╗╚██╗ ██╔╝")
    print("                    ███████╗██║   ██║██╔████╔██║██╔████╔██║███████║██████╔╝ ╚████╔╝ ")
    print("                    ╚════██║██║   ██║██║╚██╔╝██║██║╚██╔╝██║██╔══██║██╔══██╗  ╚██╔╝  ")
    print("                    ███████║╚██████╔╝██║ ╚═╝ ██║██║ ╚═╝ ██║██║  ██║██║  ██║   ██║   ")
    print("                    ╚══════╝ ╚═════╝ ╚═╝     ╚═╝╚═╝     ╚═╝╚═╝  ╚═╝╚═╝  ╚═╝   ╚═╝   ")
    print("="*80)
    print("🗂️  FOLDER-BASED PARSING COMPLETED")
    print("="*80)
    print(f"📁 INPUT DIRECTORY: {input_dir}")
    print(f"🔢 TOTAL CHUNKS: {total_chunks:,}")
    print(f"📂 FOLDERS PROCESSED: {len(all_folder_summaries)}")
    
    # Cache statistics
    total_files = total_parsed + total_cached
    if not args.no_cache and total_files > 0:
        print(f"⚡ CACHE STATISTICS:")
        print(f"   🔄 Files parsed: {total_parsed}")
        print(f"   💾 Files from cache: {total_cached}")
        print(f"   📈 Cache hit rate: {total_cached/total_files*100:.1f}%")
    elif args.no_cache:
        print(f"🚫 CACHING DISABLED - All {total_files} files were parsed")
    
    # Detailed folder breakdown
    print(f"📊 FOLDER BREAKDOWN:")
    for summary in sorted(all_folder_summaries, key=lambda x: x['chunks'], reverse=True):
        print(f"   📁 {summary['folder_key']}: {summary['chunks']:,} chunks ({summary['files_parsed'] + summary['files_cached']} files)")
        print(f"      📄 Output: {os.path.basename(summary['output_path'])}")
        if summary['images_extracted'] > 0:
            print(f"      🖼️  Images: {summary['images_extracted']}")
    
    # Show next steps
    print(f"🚀 NEXT STEPS:")
    print(f"   Run chunking and embedding for each folder:")
    for summary in all_folder_summaries:
        chunk_output = summary['output_path'].replace('parsed_', 'chunked_')
        embed_output = summary['output_path'].replace('parsed_', 'embedded_').replace('.jsonl', '.npz')
        print(f"   python pipeline/chunk.py --input \"{summary['output_path']}\" --output \"{chunk_output}\"")
        print(f"   python pipeline/embed.py --input \"{chunk_output}\" --output \"{embed_output}\"")
    
    print("="*80)
    print("✅ FOLDER-BASED PARSING COMPLETED SUCCESSFULLY!")
    print("="*80 + "\n")
    
    logging.info("Folder-based parsing completed successfully. Processed %d folders with %d total chunks.", 
                len(all_folder_summaries), total_chunks)
    return 0


if __name__ == "__main__":  # Allows `python pipeline/parse.py` execution
    sys.exit(main())


